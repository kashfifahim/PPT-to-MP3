from flask import Flask, request, render_template, send_file
from pptx import Presentation
from io import BytesIO
from gtts import gTTS
from pathlib import Path
import re
import tempfile
import os
from flask_socketio import SocketIO
import logging

app = Flask(__name__)
socketio = SocketIO(app, cors_allowed_origins="*")

# Ensure the 'uploads/' directory exists
uploads_dir = 'uploads'
if not os.path.exists(uploads_dir):
    os.makedirs(uploads_dir)

# Configure logging
logging.basicConfig(filename='app.log', level=logging.DEBUG)


@socketio.on('connect')
def test_connect():
    print('Client connected')
    try:
        logging.info('Client connected')
    except Exception as e:
        logging.error(f'Error in test_connect: {e}')


@socketio.on('disconnect')
def test_disconnect():
    print('Client disconnected')
    try:
        logging.info('Client disconnected')
    except Exception as e:
        logging.error(f'Error in test_disconnect: {e}')


# Define the clean_up_text function
def clean_up_text(input_text):
    try:
        # Add spaces after periods to separate sentences
        cleaned_text = re.sub(r'([.!?])', r'\1 ', input_text)

        # Replace multiple spaces with a single space
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

        # Capitalize the first letter of each sentence
        sentences = cleaned_text.split('. ')
        cleaned_text = '. '.join(sentence.capitalize() for sentence in sentences)

        return cleaned_text
    except Exception as e:
        logging.error(f'Error in clean_up_text: {e}')
        return input_text


def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return ''.join(text_runs).strip()
    except Exception as e:
        logging.error(f'Error in extract_text_from_pptx: {e}')
        return ""


def generate_mp3(text):
    try:
        cleaned_string = clean_up_text(text)
        tts = gTTS(cleaned_string, lang='en')
        mp3_fp = BytesIO()
        tts.write_to_fp(mp3_fp)
        return mp3_fp
    except Exception as e:
        logging.error(f'Error in generate_mp3: {e}')
        return None


def save_mp3(mp3_fp, mp3_filename):
    try:
        with open(mp3_filename, 'wb') as file:
            file.write(mp3_fp.getvalue())
    except Exception as e:
        logging.error(f'Error in save_mp3: {e}')


def process_pptx_file(pptx_file):
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_file_path = os.path.join(temp_dir, pptx_file.filename)
            pptx_file.save(temp_file_path)

            text = extract_text_from_pptx(temp_file_path)
            if text:
                mp3_fp = generate_mp3(text)
                if mp3_fp:
                    mp3_filename = os.path.join(uploads_dir, pptx_file.filename + '.mp3')
                    save_mp3(mp3_fp, mp3_filename)
                                        # Delete the processed pptx file
                    os.remove(temp_file_path)
                    return mp3_filename
    except Exception as e:
        logging.error(f'Error in process_pptx_file: {e}')
    return None


@app.route('/', methods=['GET', 'POST'])
def index():
    try:
        if request.method == 'POST':
            if 'pptx_file' in request.files:
                pptx_file = request.files['pptx_file']
                if pptx_file.filename != '':
                    # Notify the client that the conversion has started
                    socketio.emit('conversion_started', namespace='/test')

                    mp3_filename = process_pptx_file(pptx_file)

                    if mp3_filename:
                        # Notify the client that the conversion is complete
                        socketio.emit('conversion_completed', namespace='/test')

                        return send_file(mp3_filename, as_attachment=True)

        return render_template('index.html')
    except Exception as e:
        logging.error(f'Error in index route: {e}')
        return "An error occurred"


if __name__ == '__main__':
    app.run(debug=True, port=5334)